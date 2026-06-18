"""PPTX 解析器"""
import logging
from typing import List
from pptx import Presentation
from src.parser.base import BaseParser, ParseResult
logger = logging.getLogger("course_assistant.parser.pptx")

class PPTXParser(BaseParser):
    @property
    def supported_extensions(self) -> List[str]: return [".pptx", ".PPTX"]

    def parse(self, file_path):
        result = ParseResult()
        try:
            prs = Presentation(file_path); slides, notes = [], []
            for sn, slide in enumerate(prs.slides, 1):
                texts = [s.text.strip() for s in slide.shapes if s.has_text_frame and s.text.strip()]
                if texts: slides.append(f"## 幻灯片 {sn}\n" + "\n".join(texts))
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    nt = slide.notes_slide.notes_text_frame.text.strip()
                    if nt: notes.append(f"## 备注 {sn}\n{nt}")
            result.text = "\n\n".join(slides)
            if notes: result.text += "\n\n" + "\n\n".join(notes)
            result.metadata = {"slides": len(prs.slides), "has_notes": bool(notes)}
        except Exception as e: result.errors.append(f"PPTX error: {e}")
        return result
